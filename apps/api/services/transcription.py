"""
Service de transcription audio/vidéo.

Priorité : Mistral (Voxtral) si MISTRAL_API_KEY, sinon OpenAI Whisper si
OPENAI_API_KEY, sinon transcription mock. Mistral est appelé via son endpoint
REST (httpx) — pas d'endpoint OpenAI-compatible pour la transcription, mais le
format multipart suit la même convention que Whisper.
"""
import io
import logging
import tempfile
import os
from typing import Optional

import httpx

from core.config import settings

logger = logging.getLogger(__name__)

# Modèle de transcription Mistral (Voxtral). Voir docs.mistral.ai (audio).
MISTRAL_TRANSCRIBE_URL = "https://api.mistral.ai/v1/audio/transcriptions"
MISTRAL_TRANSCRIBE_MODEL = "voxtral-mini-latest"

MOCK_TRANSCRIPTION = """
Bienvenue dans ce cours sur les fondamentaux de l'intelligence artificielle.
Aujourd'hui, nous allons explorer les concepts clés qui fondent le machine learning moderne.

Commençons par définir ce qu'est l'apprentissage automatique. Il s'agit d'un ensemble de techniques
permettant à un système informatique d'apprendre à partir de données, sans être explicitement programmé.

La première notion fondamentale est celle de modèle. Un modèle est une fonction mathématique
qui prend des données en entrée et produit une prédiction en sortie.

Ensuite, l'entraînement : c'est le processus par lequel le modèle ajuste ses paramètres
pour minimiser l'erreur sur un jeu de données d'entraînement.

Nous verrons également la validation croisée, technique essentielle pour évaluer
la capacité de généralisation d'un modèle sur de nouvelles données.

Pour conclure cette introduction, retenez que le machine learning repose sur trois piliers :
les données, les algorithmes, et la puissance de calcul.
Dans les prochaines sections, nous approfondirons chacun de ces aspects.
""".strip()


def _transcribe_mistral(file_bytes: bytes, filename: str) -> Optional[str]:
    """Transcrit via Mistral Voxtral (endpoint REST). None si échec → fallback appelant."""
    try:
        resp = httpx.post(
            MISTRAL_TRANSCRIBE_URL,
            headers={"Authorization": f"Bearer {settings.MISTRAL_API_KEY}"},
            files={"file": (filename, file_bytes, "application/octet-stream")},
            data={"model": MISTRAL_TRANSCRIBE_MODEL},
            timeout=300,  # une vidéo peut être longue à transcrire
        )
        if resp.status_code >= 400:
            logger.error("Transcription Mistral %s : %s", resp.status_code, resp.text[:300])
            return None
        return resp.json().get("text")
    except Exception as e:
        logger.error("Erreur transcription Mistral: %s", e)
        return None


def _transcribe_whisper(file_bytes: bytes, filename: str) -> Optional[str]:
    """Transcrit via OpenAI Whisper. None si échec → fallback appelant."""
    try:
        from openai import OpenAI
        client = OpenAI(api_key=settings.OPENAI_API_KEY)

        ext = filename.rsplit(".", 1)[-1] if "." in filename else "mp4"
        with tempfile.NamedTemporaryFile(suffix=f".{ext}", delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            with open(tmp_path, "rb") as f:
                transcript = client.audio.transcriptions.create(
                    model="whisper-1", file=f, language="fr",
                )
            return transcript.text
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        logger.error("Erreur Whisper API: %s", e)
        return None


async def transcribe_audio(file_bytes: bytes, filename: str) -> str:
    """
    Transcrit un fichier audio/vidéo.
    Priorité Mistral (Voxtral) → Whisper → transcription mock (si rien de configuré
    ou en cas d'échec, pour ne jamais faire échouer la génération de cours).
    """
    if settings.MISTRAL_API_KEY:
        text = _transcribe_mistral(file_bytes, filename)
        if text:
            return text

    if settings.OPENAI_API_KEY:
        text = _transcribe_whisper(file_bytes, filename)
        if text:
            return text

    logger.info("Transcription indisponible (aucune clé ou échec) — transcription mock utilisée")
    return MOCK_TRANSCRIPTION


def fetch_youtube_transcript(youtube_url: str) -> Optional[str]:
    """
    Récupère les sous-titres d'une vidéo YouTube (gratuit, sans transcription LLM →
    économise des tokens). Essaie le français puis l'anglais. None si indisponible
    (sous-titres désactivés, vidéo privée, IP serveur bloquée…) → l'appelant gère.
    """
    video_id = extract_youtube_id(youtube_url) or youtube_url.strip()
    if not video_id:
        return None
    try:
        from youtube_transcript_api import YouTubeTranscriptApi

        segments = YouTubeTranscriptApi.get_transcript(video_id, languages=["fr", "en"])
        text = " ".join(s["text"] for s in segments if s.get("text"))
        return text.strip() or None
    except Exception as e:
        logger.warning("Sous-titres YouTube indisponibles (%s) : %s", video_id, e)
        return None


def extract_youtube_id(url: str) -> Optional[str]:
    """Extrait l'ID YouTube depuis une URL."""
    import re
    patterns = [
        r"(?:youtube\.com/watch\?v=|youtu\.be/|youtube\.com/embed/)([a-zA-Z0-9_-]{11})",
    ]
    for p in patterns:
        m = re.search(p, url)
        if m:
            return m.group(1)
    return None


def extract_pptx_text(file_bytes: bytes) -> str:
    """Extrait le texte d'un fichier PPTX slide par slide."""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        lines = []
        for i, slide in enumerate(prs.slides):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                lines.append(f"[Slide {i + 1}] {' | '.join(slide_texts)}")
        return "\n".join(lines) if lines else "[Aucun texte extrait des slides]"
    except ImportError:
        logger.warning("python-pptx non installé")
        return "[Extraction PPTX indisponible — python-pptx requis]"
    except Exception as e:
        logger.error("Erreur extraction PPTX: %s", e)
        return f"[Erreur extraction PPTX: {e}]"


def extract_pdf_text(file_bytes: bytes) -> str:
    """Extrait le texte d'un fichier PDF page par page."""
    try:
        import pdfplumber
        lines = []
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            for i, page in enumerate(pdf.pages):
                text = page.extract_text()
                if text and text.strip():
                    lines.append(f"[Page {i + 1}] {text.strip()}")
        return "\n".join(lines) if lines else "[Aucun texte extrait du PDF]"
    except ImportError:
        logger.warning("pdfplumber non installé")
        return "[Extraction PDF indisponible — pdfplumber requis]"
    except Exception as e:
        logger.error("Erreur extraction PDF: %s", e)
        return f"[Erreur extraction PDF: {e}]"
