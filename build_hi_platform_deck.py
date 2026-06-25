"""
Génère le deck de présentation Hi! Platform au format .pptx.

UTILISATION
-----------
• Sur claude.ai (version web) : joins ce fichier et demande
  « Exécute ce script pour générer le PowerPoint ». Le bac à sable
  contient déjà python-pptx.
• En local :  pip install python-pptx  puis  python build_hi_platform_deck.py
  → produit  hi-platform-deck.pptx  dans le dossier courant.

Design : couleurs Hi! PARIS (bleu #1A3A8F, rouge #D72638, navy #0B1D3A),
slides sombres pour ouverture / IA / clôture, claires pour le contenu.
Police : Calibri (sûre, rendue partout). Remplace FONT par "DM Sans" si
cette police est installée sur le poste de présentation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0B, 0x1D, 0x3A)
NAVY2  = RGBColor(0x07, 0x12, 0x24)
PRIMARY = RGBColor(0x1A, 0x3A, 0x8F)
PRIMARY_L = RGBColor(0x6F, 0x8E, 0xE0)
DANGER = RGBColor(0xD7, 0x26, 0x38)
DANGER_L = RGBColor(0xFF, 0x7A, 0x86)
SURFACE = RGBColor(0xF4, 0xF6, 0xFA)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
INK    = RGBColor(0x11, 0x20, 0x3A)
MUTED  = RGBColor(0x5B, 0x64, 0x79)
LINE   = RGBColor(0xE1, 0xE6, 0xF0)
DIM    = RGBColor(0x9A, 0xA3, 0xB5)
CLOUD  = RGBColor(0xC4, 0xCF, 0xE6)      # texte clair sur fond sombre
CARD_D = RGBColor(0x14, 0x2A, 0x50)      # carte sur fond sombre

FONT = "Calibri"

# ── Setup ─────────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
MX = 0.9  # marge horizontale


def new_slide(bg):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background(); r.shadow.inherit = False
    return s


def text(s, txt, x, y, w, h, size, color, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, spacing=1.05, font=FONT, italic=False):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, line in enumerate(str(txt).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = spacing
        run = p.add_run(); run.text = line
        f = run.font
        f.size = Pt(size); f.bold = bold; f.italic = italic
        f.color.rgb = color; f.name = font
    return tb


def card(s, x, y, w, h, fill, line=None):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is not None:
        shp.line.color.rgb = line; shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    try:
        shp.adjustments[0] = 0.06
    except Exception:
        pass
    return shp


def icon(s, x, y, emoji, bg, size=0.62):
    chip = card(s, x, y, size, size, bg)
    text(s, emoji, x, y, size, size, 18, WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return chip


def header(s, kicker, title, dark=False, title_size=33):
    text(s, kicker.upper(), MX, 0.72, 11, 0.4, 13, PRIMARY_L if dark else PRIMARY, bold=True)
    text(s, title, MX, 1.18, 11.5, 1.4, title_size, WHITE if dark else INK, bold=True, spacing=1.02)


def footer(s, num, dark=False):
    c = DIM if not dark else RGBColor(0x6F, 0x7C, 0x97)
    tb = s.shapes.add_textbox(Inches(MX), Inches(7.02), Inches(4), Inches(0.3))
    tf = tb.text_frame; tf.margin_left = 0; tf.margin_top = 0
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = "Hi! "; r1.font.bold = True
    r1.font.color.rgb = (PRIMARY_L if dark else PRIMARY); r1.font.size = Pt(11); r1.font.name = FONT
    r2 = p.add_run(); r2.text = "Platform"; r2.font.color.rgb = c; r2.font.size = Pt(11); r2.font.name = FONT
    text(s, num, 11.43, 7.02, 1, 0.3, 11, c, align=PP_ALIGN.RIGHT, bold=True)


def grid_cards(s, items, cols, x, y, total_w, card_h, dark=False, gap=0.35, accent_idx=None):
    cw = (total_w - gap * (cols - 1)) / cols
    fill = CARD_D if dark else WHITE
    ln = None if dark else LINE
    for k, (emoji, title, desc) in enumerate(items):
        col = k % cols; row = k // cols
        cx = x + col * (cw + gap); cy = y + row * (card_h + gap)
        ic_bg = RGBColor(0x2A, 0x46, 0x7A) if dark else RGBColor(0xE7, 0xEC, 0xF8)
        if accent_idx is not None and k == accent_idx:
            ic_bg = RGBColor(0x6A, 0x16, 0x20) if dark else RGBColor(0xFA, 0xE2, 0xE5)
        card(s, cx, cy, cw, card_h, fill, ln)
        icon(s, cx + 0.32, cy + 0.32, emoji, ic_bg)
        text(s, title, cx + 0.32, cy + 1.12, cw - 0.6, 0.5, 16, WHITE if dark else INK, bold=True)
        text(s, desc, cx + 0.32, cy + 1.62, cw - 0.6, card_h - 1.7, 11.5,
             CLOUD if dark else MUTED, spacing=1.1)


# ════════════════════════════════════════════════════════════════════════════
# 1 — TITRE
s = new_slide(NAVY)
tb = s.shapes.add_textbox(Inches(MX), Inches(0.7), Inches(8), Inches(0.6))
p = tb.text_frame.paragraphs[0]
r = p.add_run(); r.text = "Hi! "; r.font.bold = True; r.font.size = Pt(22); r.font.name = FONT; r.font.color.rgb = DANGER_L
r = p.add_run(); r.text = "Platform"; r.font.bold = True; r.font.size = Pt(22); r.font.name = FONT; r.font.color.rgb = WHITE
text(s, "PLATEFORME PÉDAGOGIQUE — HI! PARIS", MX, 2.95, 11, 0.4, 14, PRIMARY_L, bold=True)
text(s, "Le LMS nouvelle génération :\nmoderne, souverain, augmenté par l'IA.", MX, 3.45, 11.5, 2.0, 38, WHITE, bold=True, spacing=1.04)
text(s, "Une plateforme d'enseignement unifiée pour les écoles de l'Institut Polytechnique de Paris et HEC — pensée pour remplacer Moodle et les LMS externes.",
     MX, 5.35, 9.6, 1.1, 17, CLOUD, spacing=1.3)
text(s, "École Polytechnique   ·   HEC Paris   ·   Télécom Paris   ·   ENSAE   ·   Télécom SudParis   ·   ENSTA",
     MX, 6.65, 11.5, 0.4, 12, RGBColor(0x8F, 0xA0, 0xC4), bold=True)

# 2 — CONSTAT
s = new_slide(SURFACE)
header(s, "Le constat", "Les LMS actuels freinent l'enseignement")
bullets = ("Interfaces datées (« années 90 »), faible engagement\n"
           "Outils dispersés : un LMS différent par école, par cours\n"
           "Aucune mutualisation : le même cours recréé sans cesse\n"
           "Dépendance à des solutions externes, peu de contrôle des données\n"
           "Création de contenu lente et chronophage pour les enseignants")
text(s, bullets, MX, 2.7, 6.6, 3.6, 17, INK, spacing=1.7)
card(s, 8.0, 2.7, 4.43, 3.4, WHITE, LINE)
icon(s, 8.35, 3.05, "🧩", RGBColor(0xFA, 0xE2, 0xE5))
text(s, "Beaucoup d'efforts, peu de levier", 8.35, 3.9, 3.8, 0.6, 17, INK, bold=True)
text(s, "Chaque école réinvestit le même temps dans les mêmes contenus, sur des outils que personne ne maîtrise vraiment.",
     8.35, 4.5, 3.75, 1.4, 13, MUTED, spacing=1.25)
footer(s, "02")

# 3 — VISION
s = new_slide(NAVY)
header(s, "La vision", "Une seule plateforme : mutualisée, simple, maîtrisée — un « Moodle 2.0 ».", dark=True, title_size=29)
grid_cards(s, [
    ("⚡", "Simplicité radicale", "Créer un cours sans 40 boutons. Y accéder en 3 clics."),
    ("🤝", "Mutualisée", "Les écoles partagent et réutilisent les briques de cours."),
    ("🛡️", "Souveraine", "Open source, contrôlée en interne, données en Europe."),
], 3, MX, 2.85, 11.53, 2.5, dark=True)
text(s, "Inspirations UX : l'engagement de YouTube, l'interactivité de Google Colab, la certification de HFactory.",
     MX, 5.7, 11.5, 0.8, 16, CLOUD, spacing=1.3)
footer(s, "03", dark=True)

# 4 — STAT
s = new_slide(SURFACE)
header(s, "Pourquoi maintenant", "La mutualisation change l'échelle")
text(s, "×40", MX - 0.05, 2.7, 5.2, 2.6, 150, DANGER, bold=True)
text(s, "un même cours de Machine Learning peut être donné jusqu'à 40 fois, réparti dans 8 écoles.",
     MX, 5.3, 5.0, 1.4, 17, INK, spacing=1.3)
text(s, "En partageant les contenus une seule fois :", 7.2, 2.85, 5.2, 0.6, 18, INK, bold=True)
text(s, "Économies d'échelle massives sur la production de cours\n"
        "Qualité homogène pour tous les apprenants\n"
        "Les enseignants se concentrent sur la pédagogie, pas la plomberie",
     7.2, 3.6, 5.2, 2.6, 16, MUTED, spacing=1.7)
footer(s, "04")

# 5 — MODULES
s = new_slide(SURFACE)
header(s, "Architecture fonctionnelle", "Une plateforme, sept modules")
grid_cards(s, [
    ("🎬", "Hi! Tube", "Vidéothèque style YouTube, commentaires horodatés."),
    ("📖", "Hi! Course", "Cours interactifs par blocs (markdown, quiz, code)."),
    ("🎓", "Hi! MOOC", "Parcours multi-cours, prérequis et progression."),
    ("⚡", "Hi! App", "Hébergement d'apps interactives (Streamlit, Gradio)."),
    ("🛠️", "Hi! Studio", "Builder de contenu augmenté par l'IA."),
    ("🔬", "Hi! Insights", "Articles de recherche & éditoriaux."),
], 3, MX, 2.55, 11.53, 1.95, gap=0.32, accent_idx=4)
text(s, "+ Hi! Cert — certification & badges (à venir)", MX, 6.55, 8, 0.3, 12, DIM, bold=True)
footer(s, "05")

# 6 — APPRENANT
s = new_slide(SURFACE)
header(s, "Côté apprenant", "Une expérience d'apprentissage fluide")
text(s, "Catalogue riche : cours, vidéos, MOOCs, applications\n"
        "Parcours structurés avec progression et prérequis\n"
        "Quiz interactifs et suivi bloc par bloc\n"
        "Vidéos avec commentaires horodatés\n"
        "Badges & certificats à la complétion",
     MX, 2.75, 6.4, 3.5, 17, INK, spacing=1.7)
card(s, 7.9, 2.7, 4.53, 3.5, WHITE, LINE)
icon(s, 8.25, 3.05, "🧭", RGBColor(0xE7, 0xEC, 0xF8))
text(s, "« Mon Parcours »", 8.25, 3.9, 3.9, 0.5, 17, INK, bold=True)
text(s, "Chaque apprenant retrouve ses cours, sa progression et ses prochaines étapes au même endroit — en 3 clics.",
     8.25, 4.45, 3.85, 1.6, 13, MUTED, spacing=1.25)
footer(s, "06")

# 7 — STUDIO
s = new_slide(NAVY)
header(s, "Hi! Studio — le différenciateur", "Créer un cours en minutes, pas en jours", dark=True)
text(s, "Le Studio transforme les contenus existants des enseignants en ressources pédagogiques, automatiquement, via l'IA.",
     MX, 2.0, 11.3, 0.7, 16, CLOUD, spacing=1.25)
grid_cards(s, [
    ("📊", "Excel → Quiz", "Un tableur devient un QCM structuré."),
    ("🎬", "Vidéo → Cours", "Vidéo + slides → cours complet."),
    ("🃏", "Flashcards", "Cartes de révision en un clic."),
    ("🧠", "Carte mentale", "L'arborescence des concepts."),
    ("📝", "Fiche de révision", "Résumé + concepts + points clés."),
    ("❓", "FAQ auto", "Les questions des apprenants, anticipées."),
], 3, MX, 2.9, 11.53, 1.78, dark=True, gap=0.3)
footer(s, "07", dark=True)

# 8 — PIPELINE VIDEO -> COURS
s = new_slide(SURFACE)
header(s, "Zoom — pipeline Vidéo → Cours", "De la captation au cours publié")
steps = [
    ("SOURCE", "Lien YouTube", "sous-titres récupérés gratuitement", WHITE, INK),
    ("OU", "Upload vidéo", "transcription par Voxtral (Mistral)", WHITE, INK),
    ("CONTEXTE", "Slides PPTX/PDF", "texte extrait automatiquement", WHITE, INK),
    ("IA", "Cours Markdown", "structuré + quiz intégrés", NAVY, WHITE),
]
sw = 2.7; gap = 0.28; y = 3.1; x = MX
for i, (tag, title, desc, fill, fg) in enumerate(steps):
    card(s, x, y, sw, 2.0, fill, None if fill == NAVY else LINE)
    text(s, tag, x + 0.28, y + 0.3, sw - 0.5, 0.4, 11, PRIMARY_L if fill == NAVY else PRIMARY, bold=True)
    text(s, title, x + 0.28, y + 0.78, sw - 0.5, 0.5, 16, fg, bold=True)
    text(s, desc, x + 0.28, y + 1.3, sw - 0.5, 0.6, 11.5, CLOUD if fill == NAVY else MUTED, spacing=1.1)
    if i < 3:
        text(s, "→", x + sw - 0.02, y + 0.7, 0.45, 0.6, 22, PRIMARY, align=PP_ALIGN.CENTER)
    x += sw + gap
text(s, "L'enseignant choisit sa source. Les liens YouTube ne consomment aucun jeton de transcription — coût maîtrisé.",
     MX, 5.6, 11.5, 0.8, 16, INK, spacing=1.3)
footer(s, "08")

# 9 — IA SOUVERAINE
s = new_slide(NAVY)
header(s, "Une IA souveraine & maîtrisée", "Indépendante du fournisseur, par conception", dark=True)
text(s, "Couche LLM-agnostique : le modèle est un simple réglage\n"
        "Mistral en production — souveraineté européenne\n"
        "Swappable (Mistral / Claude / OpenAI) sans changer le code\n"
        "Broker de jetons à venir : crédits plafonnés, clés par utilisateur",
     MX, 2.75, 6.6, 3.4, 17, CLOUD, spacing=1.7)
card(s, 8.0, 2.7, 4.43, 3.4, CARD_D)
icon(s, 8.35, 3.05, "🇪🇺", RGBColor(0x2A, 0x46, 0x7A))
text(s, "Pourquoi ça compte", 8.35, 3.9, 3.8, 0.5, 17, WHITE, bold=True)
text(s, "Le meilleur modèle change tous les ~3 mois. En découplant l'IA du métier, la plateforme reste à jour, conforme et au meilleur coût.",
     8.35, 4.45, 3.75, 1.5, 13, CLOUD, spacing=1.25)
footer(s, "09", dark=True)

# 10 — ANALYTICS
s = new_slide(SURFACE)
header(s, "Pilotage & suivi", "Voir, comprendre, accompagner")
grid_cards(s, [
    ("📊", "Tableau de bord", "KPIs plateforme : cohortes, contenus, engagement."),
    ("⚠️", "Apprenants à risque", "Détection par inactivité & score, seuils réglables."),
    ("📁", "Exports", "Données utilisateurs & cours exportables (CSV)."),
], 3, MX, 2.7, 11.53, 2.4, accent_idx=1)
text(s, "À venir : classification automatique des élèves à suivre, et indicateurs d'apprentissage au-delà du temps passé.",
     MX, 5.55, 11.5, 0.7, 16, INK, spacing=1.3)
footer(s, "10")

# 11 — SECURITE / RGPD
s = new_slide(SURFACE)
header(s, "Confiance & conformité", "Sécurité et RGPD, nativement")
text(s, "Authentification par email institutionnel + code à usage unique\n"
        "Sessions JWT, cookies sécurisés, rotation des jetons\n"
        "Droits RGPD intégrés (accès, rectification, effacement, portabilité)\n"
        "Domaines autorisés configurables par établissement",
     MX, 2.75, 6.5, 3.4, 17, INK, spacing=1.7)
card(s, 7.95, 2.7, 4.48, 3.4, WHITE, LINE)
icon(s, 8.3, 3.05, "🔐", RGBColor(0xE7, 0xEC, 0xF8))
text(s, "Conçu pour l'institution", 8.3, 3.9, 3.9, 0.5, 17, INK, bold=True)
text(s, "Aucun mot de passe à gérer, accès réservé aux membres des écoles partenaires, données hébergées en Europe.",
     8.3, 4.45, 3.85, 1.5, 13, MUTED, spacing=1.25)
footer(s, "11")

# 12 — STACK
s = new_slide(NAVY)
header(s, "Sous le capot", "Une stack moderne et déployée", dark=True)
grid_cards(s, [
    ("🖥️", "Frontend", "Next.js 14 + Tailwind, déployé sur Vercel."),
    ("⚙️", "Backend", "FastAPI (Python), API REST typée, sur Render."),
    ("🗄️", "Données", "PostgreSQL (Supabase), migrations versionnées."),
], 3, MX, 2.85, 11.53, 2.5, dark=True)
text(s, "Architecture modulaire · CI/CD automatisée · cap sur une infra souveraine (OVH / Scaleway) en 2026.",
     MX, 5.7, 11.5, 0.8, 16, CLOUD, spacing=1.3)
footer(s, "12", dark=True)

# 13 — ROADMAP
s = new_slide(SURFACE)
header(s, "Trajectoire", "Feuille de route")
phases = [
    ("LIVRÉ", RGBColor(0x1F, 0x9D, 0x57), "Le socle",
     "Cours, vidéos, MOOCs, parcours\nStudio IA (6 générateurs)\nInsights, analytics, RGPD\nDéploiement bout-en-bout"),
    ("EN COURS / SUIVANT", PRIMARY, "Approfondir",
     "Broker LLM (coûts)\nTuteur IA (DeepTutor)\nAnalytics & classification ML"),
    ("HORIZON", DANGER, "Élargir",
     "Application mobile\nMulti-public : Bachelor, Exec Ed, grand public\nRéutilisation collaborative entre écoles"),
]
cw = (11.53 - 0.7) / 3; x = MX; y = 2.7
for tag, tc, title, body in phases:
    card(s, x, y, cw, 3.5, WHITE, LINE)
    text(s, tag, x + 0.3, y + 0.32, cw - 0.6, 0.4, 11, tc, bold=True)
    text(s, title, x + 0.3, y + 0.78, cw - 0.6, 0.5, 17, INK, bold=True)
    text(s, body, x + 0.3, y + 1.45, cw - 0.6, 1.9, 12.5, MUTED, spacing=1.45)
    x += cw + 0.35
footer(s, "13")

# 14 — DEMO
s = new_slide(NAVY)
text(s, "PLACE À LA PRATIQUE", MX, 1.5, 11, 0.4, 14, PRIMARY_L, bold=True)
text(s, "Démonstration live", MX, 2.0, 11, 1.1, 44, WHITE, bold=True)
text(s, "Connexion par email institutionnel (code à usage unique)\n"
        "Parcours d'un cours & d'un MOOC côté apprenant\n"
        "Hi! Studio : générer un quiz et des flashcards en direct\n"
        "Publier un article Insights",
     MX, 3.4, 10.5, 2.4, 18, CLOUD, spacing=1.7)
text(s, "plateform-poc-hiparis.vercel.app", MX, 6.1, 10, 0.6, 18, PRIMARY_L, bold=True)
footer(s, "14", dark=True)

# 15 — CLOTURE
s = new_slide(NAVY)
tb = s.shapes.add_textbox(Inches(MX), Inches(2.2), Inches(8), Inches(0.6))
p = tb.text_frame.paragraphs[0]
r = p.add_run(); r.text = "Hi! "; r.font.bold = True; r.font.size = Pt(22); r.font.name = FONT; r.font.color.rgb = DANGER_L
r = p.add_run(); r.text = "Platform"; r.font.bold = True; r.font.size = Pt(22); r.font.name = FONT; r.font.color.rgb = WHITE
text(s, "Apprendre, créer et partager —\nà l'échelle de Hi! PARIS.", MX, 3.0, 11.5, 1.8, 40, WHITE, bold=True, spacing=1.05)
text(s, "Moderne pour les étudiants. Simple pour les enseignants. Souverain pour l'institution.",
     MX, 4.9, 10.5, 0.8, 18, CLOUD, spacing=1.3)
text(s, "Merci — Questions ?", MX, 6.4, 8, 0.5, 15, PRIMARY_L, bold=True)
footer(s, "15", dark=True)

# ── Sauvegarde ────────────────────────────────────────────────────────────────
OUT = "hi-platform-deck.pptx"
prs.save(OUT)
print(f"✅ Présentation générée : {OUT}  ({len(prs.slides)} slides)")
